from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import io

import pandas as pd

def parse_pdf(file_bytes):
    pdf_stream = io.BytesIO(file_bytes)   # 👈 wrap bytes
    reader = PdfReader(pdf_stream)

    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""

    return text

def parse_docx(file):
    doc = Document(file)
    return "\n".join([p.text for p in doc.paragraphs])

def parse_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_excel(file):
    df = pd.read_excel(file)
    return df.to_string()
